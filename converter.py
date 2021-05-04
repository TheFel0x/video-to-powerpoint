# bruh I dont even have powerpoint to test this 
from pptx import Presentation
X = Presentation()

# Title Card
Layout = X.slide_layouts[0] 
first_slide = X.slides.add_slide(Layout)
first_slide.shapes.title.text = "Enjoy."    # should put something more interesting here
first_slide.placeholders[1].text = "https://github.com/thefel0x"

# add images
from pptx.util import Inches
import os

imagelist = os.listdir('./cache')   # should use tempfile or similar later
imagelist.sort()

for f in imagelist:
    print(f) # should add percentage
    layout = X.slide_layouts[5]
    newslide = X.slides.add_slide(layout)
    newslide.shapes.title.text = f
    # textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1))
    # textframe = textbox.text_frame
    # paragraph = textframe.add_paragraph()
    # paragraph.text = f    
    #add pic
    # TODO: calculate width / height depending on which one is greater with image dimensions
    pic = newslide.shapes.add_picture("./cache/"+f, 0, 0,width=Inches(9), height=Inches(5))


# TODO:
#   - add images
#   - use tempfile
#   - move extracter and converter in to 1 file
#   - option to change frame density (use every 2nd/3rd/n frame)
#   - reduce filesize somehow / make more performant (compression?)
#   - take args
#   - drink more water 

X.save("./export/export.pptx")
