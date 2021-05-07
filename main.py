import os
import cv2 # opencv-python
import tempfile
from pptx import Presentation # python-pptx
from pptx.util import Inches

def main():
    # TODO: take args
    input_file = input('Input file location: ')
    output_dir = input('Export directory: ')        
    output_name = input('Output file name (.pptx): ')
    
    # create temp dir
    temp = tempfile.TemporaryDirectory()
    cache_dir = temp.name

    # Extract frames from video
    # TODO: option to lower resolution
    print('# extracting frames')
    vidcap = cv2.VideoCapture(input_file)
    success,image = vidcap.read()
    count = 0
    while success:
        cv2.imwrite(os.path.join(cache_dir,'frame%09d.jpg' % count), image) # TODO: decide amount of leading zeros with file count 
        success,image = vidcap.read()
        if success:
            print("%09d"%count)
        count += 1
    print('Extracting frames done.')
    
    # Create PowerPoint
    print('# creating ppp')
    ppp = Presentation()
    layout = ppp.slide_layouts[0]
    
    # Add frames
    frames = os.listdir(cache_dir)
    frames.sort()

    for f in frames:
        print('adding '+f) # should add percentage
        slide = ppp.slides.add_slide(layout)
        slide.shapes.title.text = f
        # TODO: dimension calculations
        pic = slide.shapes.add_picture(os.path.join(cache_dir,f), 0, 0,width=Inches(9), height=Inches(5))
    ppp.save(os.path.join(output_dir,output_name))
    print('done. saved at:\n'+os.path.join(output_dir,output_name))

    # delete temp dir
    temp.cleanup()

if __name__ == '__main__':
    main()
